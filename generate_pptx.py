import sys
import os

# Check and install python-pptx if it is not present
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
except ImportError:
    print("python-pptx is not installed. Installing it now...")
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    
    # Custom color palette (matching Tulika Tours Brand Identity)
    forest_green = RGBColor(45, 90, 39)
    charcoal_dark = RGBColor(15, 23, 42)
    warm_white = RGBColor(250, 250, 249)
    accent_gold = RGBColor(245, 158, 11)
    slate_gray = RGBColor(100, 116, 139)
    
    # Slide 1: Title Slide (Dark Theme for Wow Effect)
    slide_layout = prs.slide_layouts[6] # Blank slide
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = charcoal_dark
    
    # Title Text Frame
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = "TULIKA TOURS & TRAVELS"
    p1.font.size = Pt(40)
    p1.font.bold = True
    p1.font.color.rgb = accent_gold
    p1.font.name = "Arial"
    
    p2 = tf.add_paragraph()
    p2.text = "Data-Driven ML & Analytics Dashboard"
    p2.font.size = Pt(24)
    p2.font.bold = True
    p2.font.color.rgb = warm_white
    p2.font.name = "Arial"
    p2.space_before = Pt(10)
    
    p3 = tf.add_paragraph()
    p3.text = "Project Presentation & Architecture Review"
    p3.font.size = Pt(14)
    p3.font.color.rgb = slate_gray
    p3.font.name = "Arial"
    p3.space_before = Pt(30)
    
    # Helper to add standard content slides
    def add_content_slide(title_text):
        slide = prs.slides.add_slide(slide_layout)
        # Light background
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = warm_white
        
        # Header Box
        headerBox = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(8.5), Inches(1))
        tf_h = headerBox.text_frame
        tf_h.word_wrap = True
        p_h = tf_h.paragraphs[0]
        p_h.text = title_text
        p_h.font.size = Pt(28)
        p_h.font.bold = True
        p_h.font.color.rgb = forest_green
        p_h.font.name = "Arial"
        
        # Content Box
        contentBox = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(8.5), Inches(5))
        tf_c = contentBox.text_frame
        tf_c.word_wrap = True
        return tf_c
    
    # Slide 2: Project Overview
    tf = add_content_slide("Project Overview & Objectives")
    
    p = tf.paragraphs[0]
    p.text = "Goal:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = charcoal_dark
    
    p_body = tf.add_paragraph()
    p_body.text = "To empower Tulika Tours with predictive capabilities and data insights, helping them understand customer booking patterns, travel preferences, and customer demographics."
    p_body.font.size = Pt(14)
    p_body.font.color.rgb = charcoal_dark
    p_body.space_after = Pt(20)
    
    p_obj = tf.add_paragraph()
    p_obj.text = "Key Objectives:"
    p_obj.font.size = Pt(18)
    p_obj.font.bold = True
    p_obj.font.color.rgb = charcoal_dark
    
    bullets = [
        "Create an interactive system to visualize traveler attributes and histories.",
        "Apply Machine Learning (K-Means Clustering) to segment customers dynamically.",
        "Perform Sentiment Analysis on reviews to measure customer satisfaction.",
        "Forecast monthly booking demand to optimize fleet and resources."
    ]
    for b in bullets:
        p_b = tf.add_paragraph()
        p_b.text = "• " + b
        p_b.font.size = Pt(14)
        p_b.font.color.rgb = charcoal_dark
        p_b.space_before = Pt(5)
        
    # Slide 3: Technical Architecture
    tf = add_content_slide("System Architecture & Tech Stack")
    
    tech_stacks = [
        ("Frontend Application", "React.js, React Router v6, Chart visualization using Recharts, and styled using modern CSS variables with premium dark/light themes."),
        ("Backend REST API", "Flask, Flask-SQLAlchemy (relational database storage), JWT-Extended for token authentication, and Flask-Mail for notification emails."),
        ("Data Science & ML Core", "Python, Pandas, NumPy, Scikit-learn (StandardScaler, KMeans), Matplotlib, and Seaborn for report generation.")
    ]
    
    for title, desc in tech_stacks:
        p_title = tf.add_paragraph()
        p_title.text = title
        p_title.font.size = Pt(16)
        p_title.font.bold = True
        p_title.font.color.rgb = forest_green
        p_title.space_before = Pt(12)
        
        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(13)
        p_desc.font.color.rgb = charcoal_dark
        p_desc.space_before = Pt(2)
        p_desc.space_after = Pt(5)

    # Slide 4: Machine Learning Implementation
    tf = add_content_slide("Machine Learning Implementation")
    
    p_intro = tf.paragraphs[0]
    p_intro.text = "1. Customer Segmentation (Unsupervised Learning)"
    p_intro.font.size = Pt(16)
    p_intro.font.bold = True
    p_intro.font.color.rgb = forest_green
    
    p_seg = tf.add_paragraph()
    p_seg.text = "• Algortihm: K-Means Clustering.\n" \
                 "• Features: Package price and destination popularity index.\n" \
                 "• Customer Classes: Classified into Budget Travelers (Green), Medium Clients (Purple), and Luxury Clients (Yellow)."
    p_seg.font.size = Pt(13)
    p_seg.font.color.rgb = charcoal_dark
    p_seg.space_before = Pt(4)
    p_seg.space_after = Pt(15)
    
    p_sent_title = tf.add_paragraph()
    p_sent_title.text = "2. Sentiment Classification"
    p_sent_title.font.size = Pt(16)
    p_sent_title.font.bold = True
    p_sent_title.font.color.rgb = forest_green
    
    p_sent = tf.add_paragraph()
    p_sent.text = "• Method: Rating-based sentiment mapping.\n" \
                 "• Sentiment Tiers: Positive (ratings >= 4), Neutral (rating = 3), and Negative (ratings < 3).\n" \
                 "• Outcome: Visualizes service quality and tourist experience feedback directly on the dashboard."
    p_sent.font.size = Pt(13)
    p_sent.font.color.rgb = charcoal_dark
    p_sent.space_before = Pt(4)

    # Slide 5: Key Visualizations & Analytics
    tf = add_content_slide("Key Visualizations & Analytics")
    
    viz_list = [
        ("Traveler Preferences", "Shows distribution of category bookings (Adventure, Beach, Culture, Wildlife, City) to determine what tours to advertise."),
        ("Customer Segmentation Plot", "Scatter plot highlighting the distinct clusters of traveler groups based on their spending power vs destination appeal."),
        ("Monthly Booking Volume", "Line chart showing historical demand trends, critical for scheduling vehicle availability and driver assignments."),
        ("Destination Popularity vs Satisfaction", "Dual-axis visualization contrasting the total booking counts with actual average star ratings.")
    ]
    
    for title, desc in viz_list:
        p_title = tf.add_paragraph()
        p_title.text = "• " + title + ": "
        p_title.font.size = Pt(14)
        p_title.font.bold = True
        p_title.font.color.rgb = charcoal_dark
        p_title.space_before = Pt(8)
        
        # Append description to the same paragraph
        run = p_title.add_run()
        run.text = desc
        run.font.bold = False
        run.font.size = Pt(13)
        run.font.color.rgb = slate_gray

    # Slide 6: Outcomes & Value
    tf = add_content_slide("Project Outcomes & Business Value")
    
    outcomes = [
        ("Secure Authentication & Management", "Includes password hashing, JWT authorization, registration workflows, and automated email alerts on login events."),
        ("Actionable Marketing Insights", "Allows marketing teams to design targeted packages for budget-conscious versus luxury-seeking tourists."),
        ("Enhanced Operations Control", "Provides operational visibility into monthly seasonal demands to forecast inventory, tours, and staff bookings.")
    ]
    
    for title, desc in outcomes:
        p_title = tf.add_paragraph()
        p_title.text = title
        p_title.font.size = Pt(15)
        p_title.font.bold = True
        p_title.font.color.rgb = forest_green
        p_title.space_before = Pt(10)
        
        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(13)
        p_desc.font.color.rgb = charcoal_dark
        p_desc.space_before = Pt(2)

    # Save presentation
    output_filename = "Tulika_Tours_Project_Presentation.pptx"
    prs.save(output_filename)
    print(f"Presentation saved successfully as: {output_filename}")

if __name__ == "__main__":
    create_presentation()
